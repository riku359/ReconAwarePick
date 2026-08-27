#!/usr/bin/env python3
"""Fig. 1: the pipeline, drawn with python-pptx as one slide.

Two blocks. The top block is one round of the feedback loop; no reconstruction is drawn
there, because the loop does not run one. The bottom block picks the full micrograph set
with the checkpoint the loop delivers and takes it through the same contamination mask
and 2D class selection to reconstruction.

The top block is laid out as a serpentine. Row 1 runs left to right (pick, mask,
surviving particles), turns at the right edge, and row 2 runs right to left (2D
classification, CryoSift, kept classes, pseudo-labels, fine-tune). Putting fine-tune
directly under CryoTransformer is what lets the returned checkpoint be one vertical
dashed line rather than a path around the block.

There is one set of class-average tiles, the coloured pair after CryoSift, because two
sets do not fit the width.

The bottom block is a single row of smaller pills. CryoTransformer's frame is orange
there, which is how the figure says it runs with the checkpoint the loop delivered. The
arrow from CryoSift to reconstruction carries raw particles and is labelled stack S,
matching the supplementary protocol figure.

The symbols are the ones Sec. 3 of the paper uses: M, C_n, M_i, C'_n, S_n, T_n, theta_n,
theta_{n+1}.

The box labels live in this script, not in the .pptx, so the deck is an output and not a
source. The photographic assets are not committed: point `--assets` at a directory
holding

    q_raw.jpg  q_bbox.jpg  q_masked.jpg      the three micrograph thumbnails
    patches/patch_<nnn>.jpg                  raw particle crops
    plain/cls_k0.png ... cls_d1.png          class averages with no burnt-in frame
    map3d_10081_gt_transparent.png           the 3D volume, transparent background

Runs standalone once those assets are in place: needs python-pptx and nothing else.
Turning the deck into the cropped PDF the paper carries takes LibreOffice and pypdf;
results/figures/pipeline_overview/README.md has that command, and the crop box.

    python build_pipeline_fig.py --assets <dir> [--out <deck.pptx>]
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "lib"))
import figure_paths                                    # noqa: E402

INK = RGBColor(0x26, 0x32, 0x3A)
GRAY = RGBColor(0x3D, 0x47, 0x4E)
LGRAY = RGBColor(0x8D, 0x95, 0x9C)
ORANGE = RGBColor(0xC2, 0x62, 0x0E)
GREEN = RGBColor(0x2C, 0xA0, 0x2C)
RED = RGBColor(0xD6, 0x27, 0x28)
PAGE = RGBColor(0x9F, 0xB4, 0xBF)
FRAME = RGBColor(0x8A, 0x8A, 0x8A)

FONT = "Calibri"
FONT_MATH = "Cambria"
# The paper's text width is 6.875 in. The figure is a 13.333 in wide PDF placed at
# \linewidth, so it lands at about 0.52 scale on the page. Matching the 10 pt body text
# therefore needs about 19 pt in the deck.
F = 19          # labels
F_PILL = 16     # boxes (the same size in both blocks)
F_PILL3 = 16    # boxes of the bottom block
F_SMALL = 16    # auxiliary labels and block tags

# --- vertical positions --------------------------------------------------
# The row spacing and heights are compressed one step against a plain layout, to save
# page height in print.
CY1 = 0.94           # centre of row 1
LBL1 = CY1 + 0.48    # labels of row 1
CY2 = 2.68           # centre of row 2
GB_DY = 0.06         # good/bad labels, dropped below the tile block
LBL2_DY = 0.38       # row 2 labels, dropped below the tile block
LOOP_TOP, LOOP_BOT = 0.34, 4.00
CKPT_Y = 4.22        # the lane the checkpoint travels along
FULL_TOP, FULL_BOT = 4.44, 6.02
CY3 = 5.10           # centre of row 3
LBL3 = CY3 + 0.50    # labels of row 3

# --- horizontal positions ------------------------------------------------
IMG = 0.84          # micrograph thumbnail
STEP = 0.05         # offset between the sheets of a micrograph stack
SP = 0.38           # a surviving-particle patch
SPG = 0.07
TILE, TGX, TGY = 0.40, 0.18, 0.17   # class-average tiles
PSTEP = 0.045       # offset of the raw particles stacked behind a class average

X0 = 0.44           # left edge (where a label's left edge clears the block frame)
RIGHT_END = 12.88   # right edge of a row
WRAP_X = 13.12      # the vertical line the serpentine turns on

W_CT, W_MC = 1.30, 1.30             # pills of the top block, sized for two lines at 16 pt
W_2D, W_CS, W_FT = 1.52, 1.10, 1.15
W_CT3, W_MC3, W_2D3, W_CS3, W_RC3 = 1.30, 1.30, 1.52, 1.10, 1.40  # bottom block, same widths
MAPH = 0.92
MAPW = MAPH * 665 / 756

grid_w = 3 * TILE + 2 * TGX
grid_h = 2 * TILE + TGY
tiles_w = grid_w + 3 * PSTEP          # the stack behind extends it to the right
surv_w = 3 * SP + 2 * SPG
stack_w = IMG + 2 * STEP

# row 1: the even spacing is what is left once the fixed widths are taken out
row1_fixed = 3 * stack_w + W_CT + W_MC + surv_w
SEP1 = (RIGHT_END - X0 - row1_fixed) / 5
AR1 = 0.44

CLS = ["cls_k0", "cls_k1", "cls_d0", "cls_k2", "cls_k3", "cls_d1"]
BACK = [("000", "002", "004"), ("006", "008", "011"), ("014", "015", "017"),
        ("019", "020", "023"), ("024", "026", "028"), ("031", "033", "035")]

SURV = ["001", "005", "010", "012", "013", "021"]

slide = None          # set by build(); the drawing helpers below place shapes on it
ASSETS = None         # set by build(); the directory holding the photographic assets


def asset(*parts) -> str:
    path = ASSETS.joinpath(*parts)
    if not path.is_file():
        raise SystemExit(f"missing asset {path}; --assets must hold the layout in the "
                         f"module docstring")
    return str(path)


def add_runs(p, parts, size=F, color=GRAY, bold=False, italic=False):
    """One paragraph from (text, kind) pairs; kind 'i' is italic, 'sub' is a subscript."""
    for text, kind in parts:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.italic = italic or kind in ("i", "sub")
        r.font.name = FONT_MATH if kind in ("i", "sub") else FONT
        if kind == "sub":
            r._r.get_or_add_rPr().set("baseline", "-25000")


def label(cx, parts, y, size=F, color=GRAY, w=3.4, align=PP_ALIGN.CENTER, bold=False):
    """A centred label. A list of lists in `parts` is one line per element."""
    lines = parts if parts and isinstance(parts[0], list) else [parts]
    tb = slide.shapes.add_textbox(Inches(cx - w / 2), Inches(y), Inches(w),
                                  Inches(0.30 * len(lines) + 0.14))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = 0.94
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        add_runs(p, ln, size=size, color=color, bold=bold)
    return tb


def wlabel(x, y, parts, size=F_SMALL, color=GRAY, italic=True, w=4.5):
    """A left-aligned one-line label, for the connecting lanes and the block tags."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.30))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    add_runs(p, parts, size=size, color=color, italic=italic)
    return tb


def tag(x, y, parts, w):
    """A tag just above a block frame's top edge, so the text never sits on the line."""
    return wlabel(x, y - 0.30, parts, size=F_SMALL, color=GRAY, italic=True, w=w)


def pill(x, w, text, cy, fill="FFFFFF", line=INK, tcolor=INK, h=0.62, size=F_PILL):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(cy - h / 2),
                                Inches(w), Inches(h))
    sp.adjustments[0] = 0.30
    sp.fill.solid()
    sp.fill.fore_color.rgb = RGBColor.from_string(fill)
    sp.line.color.rgb = line
    sp.line.width = Pt(1.75)
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 0.88
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.color.rgb = tcolor
        r.font.name = FONT
    return sp


def frame(x0, y0, x1, y1):
    """The dashed rounded rectangle around a block. No fill, and drawn first so it sits
    behind everything the block holds."""
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x0), Inches(y0),
                                Inches(x1 - x0), Inches(y1 - y0))
    sp.adjustments[0] = 0.03
    sp.fill.background()
    sp.line.color.rgb = FRAME
    sp.line.width = Pt(1.1)
    sp.shadow.inherit = False
    el = sp.line._get_or_add_ln()
    el.append(el.makeelement(qn("a:prstDash"), {"val": "dash"}))
    return sp


def line_seg(x0, y0, x1, y1, color=GRAY, dashed=False, head=False, width=1.6):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x0), Inches(y0),
                                    Inches(x1), Inches(y1))
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    el = ln.line._get_or_add_ln()
    if dashed:
        el.append(el.makeelement(qn("a:prstDash"), {"val": "dash"}))
    if head:
        el.append(el.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med",
                                                   "len": "med"}))
    return ln


def arrow_r(x, cy, sep, ar=AR1, color=GRAY):
    """A rightward arrow. `x` is the previous element's right edge; returns the next
    element's left edge."""
    g = (sep - ar) / 2
    line_seg(x + g, cy, x + g + ar, cy, head=True, color=color)
    return x + sep


def arrow_l(x_from, x_to, cy, color=GRAY):
    """A leftward arrow, from the previous element's left edge to the next one's right."""
    line_seg(x_from - 0.05, cy, x_to + 0.05, cy, head=True, color=color)


def micrograph_stack(path, x, cy, sheets=3):
    """A micrograph drawn as `sheets` overlapping sheets. Returns the right edge, which
    is fixed at three sheets' worth so the row spacing does not move with the count."""
    for k in range(sheets - 1, -1, -1):
        pic = slide.shapes.add_picture(path, Inches(x + STEP * k),
                                       Inches(cy - IMG / 2 - STEP * k * 0.8),
                                       width=Inches(IMG), height=Inches(IMG))
        pic.line.color.rgb = PAGE
        pic.line.width = Pt(1.0)
    return x + stack_w


def patch(name):
    return asset("patches", f"patch_{name}.jpg")


def particle(path, x, y, s, color=GREEN, width=1.5):
    pic = slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(s),
                                   height=Inches(s))
    pic.line.color.rgb = color
    pic.line.width = Pt(width)
    return pic


def particle_row(names, x, cy, s, g, color=GREEN, width=1.5):
    for i, n in enumerate(names):
        particle(patch(n), x + i * (s + g), cy - s / 2, s, color=color, width=width)
    return x + len(names) * s + (len(names) - 1) * g


def halo(x, y, w, h):
    """A white underlay, laid under each front patch so that the frames of overlapping
    patches never touch."""
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def particle_stack(names, x, cy, s, d=0.12, color=GREEN, width=1.5, pad=0.04):
    """Particles overlapped on a diagonal, as the supplementary protocol figure carries
    them. `names[0]` is in front, at the lower left.

    The patches behind step up and to the right, so each front patch takes a white
    margin of `pad` above and to its right and the frames do not touch. Only those two
    sides take it: the arrow passes along the left and the bottom.
    """
    h = s + (len(names) - 1) * d * 0.8
    y0 = cy + h / 2 - s          # top edge of the front patch
    for k in range(len(names) - 1, -1, -1):
        px, py = x + k * d, y0 - k * d * 0.8
        if k < len(names) - 1:                 # the rearmost patch needs no underlay
            halo(px, py - pad, s + pad, s + pad)
        particle(patch(names[k]), px, py, s, color=color, width=width)
    return x + s + (len(names) - 1) * d


def class_tiles(x0, cy, colors, widths):
    """A 3x2 grid of class averages, each with three raw particles stacked behind it."""
    y0 = cy - grid_h / 2
    for i, name in enumerate(CLS):
        row, col = divmod(i, 3)
        tx = x0 + col * (TILE + TGX)
        ty = y0 + row * (TILE + TGY)
        for k, pn in enumerate(reversed(BACK[i])):        # back to front
            d = PSTEP * (3 - k)
            particle(patch(pn), tx + d, ty - d * 0.8, TILE, color=PAGE, width=1.0)
        pic = slide.shapes.add_picture(asset("plain", f"{name}.png"),
                                       Inches(tx), Inches(ty),
                                       width=Inches(TILE), height=Inches(TILE))
        pic.line.color.rgb = colors[i]
        pic.line.width = Pt(widths[i])
    return y0 + grid_h


def build(assets: Path, out_path: Path):
    global slide, ASSETS
    ASSETS = assets

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)
    slide._element.set("showMasterSp", "0")

    # --- block frames, drawn first so they sit behind ---------------------
    frame(0.15, LOOP_TOP, 13.22, LOOP_BOT)
    frame(0.15, FULL_TOP, 13.22, FULL_BOT)

    # --- row 1 (left to right) -------------------------------------------
    # micrographs -> CryoTransformer -> picks -> MicrographCleaner
    #             -> contamination masks -> surviving particles
    x = X0
    e = micrograph_stack(asset("q_raw.jpg"), x, CY1)
    label(x + IMG / 2 + STEP, [[("CryoPPP 300", "")], [("micrographs ", ""), ("M", "i")]],
          LBL1, size=F_SMALL, w=2.2)
    x = arrow_r(e, CY1, SEP1)

    ct_x = x
    pill(ct_x, W_CT, "Cryo\nTransformer", CY1)
    ct_cx = ct_x + W_CT / 2
    x = arrow_r(x + W_CT, CY1, SEP1)

    e = micrograph_stack(asset("q_bbox.jpg"), x, CY1)
    label(x + IMG / 2 + STEP, [("picks ", ""), ("C", "i"), ("n", "sub")], LBL1)
    x = arrow_r(e, CY1, SEP1)

    pill(x, W_MC, "Micrograph\nCleaner", CY1)
    x = arrow_r(x + W_MC, CY1, SEP1)

    e = micrograph_stack(asset("q_masked.jpg"), x, CY1)
    label(x + IMG / 2 + STEP,
          [[("contamination", "")], [("masks ", ""), ("M", "i"), ("i", "sub")]],
          LBL1, w=2.2)
    x = arrow_r(e, CY1, SEP1)

    surv_end = particle_row(SURV[:3], x, CY1 - (SP + SPG) / 2, SP, SPG)
    particle_row(SURV[3:], x, CY1 + (SP + SPG) / 2, SP, SPG)
    surv_cx = x + surv_w / 2
    label(surv_cx,
          [[("surviving", "")], [("particles ", ""), ("C", "i"), ("′", "i"), ("n", "sub")]],
          LBL1, w=2.4)

    # --- the turn (right edge of row 1 to right edge of row 2) ------------
    p2d_x = RIGHT_END - W_2D
    line_seg(surv_end + 0.05, CY1, WRAP_X, CY1)
    line_seg(WRAP_X, CY1, WRAP_X, CY2)
    line_seg(WRAP_X, CY2, RIGHT_END + 0.03, CY2, head=True)

    # --- row 2 (right to left) -------------------------------------------
    # 2D classification -> CryoSift -> kept/dropped classes -> pseudo-labels -> fine-tune
    pill(p2d_x, W_2D, "2D\nclassification", CY2)

    teach_s, teach_g = 0.34, 0.07
    teach_w = 3 * teach_s + 2 * teach_g

    ft_x = ct_cx - W_FT / 2
    SEP2 = 0.90                     # tiles <-> CryoSift <-> 2D classification

    cs_x = p2d_x - SEP2 - W_CS
    arrow_l(p2d_x, cs_x + W_CS, CY2)
    pill(cs_x, W_CS, "CryoSift", CY2)

    tls_x = cs_x - SEP2 - tiles_w
    arrow_l(cs_x, tls_x + tiles_w, CY2)
    kept_colors = [GREEN if n.startswith("cls_k") else RED for n in CLS]
    tiles_bottom = class_tiles(tls_x, CY2, kept_colors, [3.0] * 6)
    label(tls_x + TILE + TGX / 2, [("good", "")], tiles_bottom + GB_DY,
          size=F_SMALL, color=GREEN, w=1.2, bold=True)
    label(tls_x + 2 * (TILE + TGX) + TILE / 2, [("bad", "")], tiles_bottom + GB_DY,
          size=F_SMALL, color=RED, w=1.2, bold=True)
    label(tls_x + grid_w / 2, [("kept classes → stack ", ""), ("S", "i"), ("n", "sub")],
          tiles_bottom + LBL2_DY, w=3.4)

    # Row 2 has width to spare, so the pseudo-labels stay in a row; only the bottom
    # block's stack S is overlapped on a diagonal.
    teach_cx = (tls_x + ft_x + W_FT) / 2
    teach_x = teach_cx - teach_w / 2
    arrow_l(tls_x, teach_x + teach_w, CY2, color=GREEN)
    particle_row(["022", "027", "030"], teach_x, CY2, teach_s, teach_g)
    label(teach_cx, [("pseudo-labels ", ""), ("T", "i"), ("n", "sub")],
          CY2 + teach_s / 2 + 0.07, size=F_SMALL, w=2.4)

    arrow_l(teach_x, ft_x + W_FT, CY2, color=GREEN)
    pill(ft_x, W_FT, "fine-tune", CY2, fill="FDF3E7", line=ORANGE, tcolor=ORANGE)

    # theta_{n+1}: from the top edge of fine-tune to the bottom edge of CryoTransformer,
    # dashed and orange. fine-tune sits directly below, so this is a single line.
    line_seg(ct_cx, CY2 - 0.31 - 0.03, ct_cx, CY1 + 0.31 + 0.03, color=ORANGE,
             dashed=True, head=True)
    wlabel(ct_cx + 0.14, (CY1 + CY2) / 2 - 0.02,
           [("θ", "i"), ("n+1", "sub"), ("  next round", "")], color=ORANGE)

    tag(0.35, LOOP_TOP, [("the feedback loop, round ", ""), ("n", "i")], w=2.50)

    # --- the lane the delivered checkpoint travels along ------------------
    ct3_x = X0 + stack_w + 0.38
    ct3_cx = ct3_x + W_CT3 / 2
    line_seg(ct_cx, CY2 + 0.31 + 0.03, ct_cx, CKPT_Y, color=ORANGE)
    line_seg(ct_cx, CKPT_Y, ct3_cx, CKPT_Y, color=ORANGE)
    line_seg(ct3_cx, CKPT_Y, ct3_cx, CY3 - 0.31 - 0.03, color=ORANGE, head=True)
    wlabel(ct_cx + 0.16, LOOP_BOT + 0.06,
           [("the checkpoint the loop delivers", "")], color=ORANGE)

    tag(0.35, FULL_TOP, [("3D reconstruction", "")], w=1.60)

    # --- row 3 (left to right) -------------------------------------------
    # full micrograph set -> CryoTransformer -> MicrographCleaner -> 2D classification
    #                     -> CryoSift -> stack S -> reconstruction -> 3D volume
    SEP3 = 0.57
    AR3 = 0.30

    x = X0
    e = micrograph_stack(asset("q_raw.jpg"), x, CY3, sheets=4)
    label(x + IMG / 2 + STEP, [("full set ", ""), ("M", "i")], LBL3, size=F_SMALL, w=2.0)
    x = arrow_r(e, CY3, SEP3, ar=AR3)

    pill(x, W_CT3, "Cryo\nTransformer", CY3, line=ORANGE, size=F_PILL3)
    x = arrow_r(x + W_CT3, CY3, SEP3, ar=AR3)

    pill(x, W_MC3, "Micrograph\nCleaner", CY3, size=F_PILL3)
    x = arrow_r(x + W_MC3, CY3, SEP3, ar=AR3)

    pill(x, W_2D3, "2D\nclassification", CY3, size=F_PILL3)
    x = arrow_r(x + W_2D3, CY3, SEP3, ar=AR3)

    pill(x, W_CS3, "CryoSift", CY3, size=F_PILL3)
    x = x + W_CS3

    # CryoSift -> reconstruction: the arrow carries the particles of stack S, as the
    # supplementary protocol figure does.
    carry_s, carry_d = 0.28, 0.14
    carry_w = carry_s + 2 * carry_d
    SEPC = carry_w + 2 * 0.26
    line_seg(x + 0.05, CY3, x + SEPC - 0.05, CY3, head=True)
    carry_x = x + (SEPC - carry_w) / 2
    particle_stack(["012", "013", "021"], carry_x, CY3, carry_s, carry_d, width=1.2)
    label(carry_x + carry_w / 2, [("stack ", ""), ("S", "i")],
          CY3 + (carry_s + 2 * carry_d * 0.8) / 2 + 0.05, size=F_SMALL, w=1.6)
    x = x + SEPC

    pill(x, W_RC3, "reconstruction", CY3, fill="F3F4F5", line=LGRAY,
         tcolor=RGBColor(0x5A, 0x62, 0x68), size=F_PILL3)
    x = arrow_r(x + W_RC3, CY3, SEP3, ar=AR3)

    slide.shapes.add_picture(asset("map3d_10081_gt_transparent.png"), Inches(x),
                             Inches(CY3 - MAPH / 2), height=Inches(MAPH))
    label(x + MAPW / 2, [("3D volume", "")], LBL3, size=F_SMALL, w=1.8)
    row3_end = x + MAPW

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"row1 {X0}..{surv_end:.2f} (SEP1 {SEP1:.2f})  "
          f"row2 ft {ft_x:.2f} teach_cx {teach_cx:.2f} tiles {tls_x:.2f} cs {cs_x:.2f} "
          f"2d {p2d_x:.2f} (SEP2 {SEP2:.2f})  row3 end {row3_end:.2f}")
    print("saved", out_path)


def main():
    parser = argparse.ArgumentParser(
        description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("--assets", type=Path, required=True,
                        help="directory holding the photographic assets listed above")
    parser.add_argument("--out", type=Path, default=None,
                        help="output deck (default "
                             "$RAPICK_FIGURES_OUT/pipeline_overview.pptx)")
    args = parser.parse_args()
    out = args.out or figure_paths.figures_out() / "pipeline_overview.pptx"
    build(args.assets.expanduser(), Path(out))


if __name__ == "__main__":
    main()
