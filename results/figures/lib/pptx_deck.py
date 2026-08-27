"""The deck helpers the two hand-laid overlay figures share.

Fig. 2 and Fig. 6 are both bare micrograph panels with a description set underneath.
The renderer that produced the panels burns a black title bar into every overlay, which
is unreadable at column width and wastes vertical space, so `prepare_overlay_panels.py`
strips it and these decks set the description as text under the image it describes.

Each deck is one slide, sized so that the exported PDF lands at about half scale on the
page. That is what keeps 18 pt here reading as roughly 9 pt in print. A figure set
narrower than its slide implies carries the extra factor in its own font sizes.

Export and crop as for the pipeline figure:

    soffice --headless --convert-to pdf --outdir <out> <deck>.pptx
"""
from __future__ import annotations

from pathlib import Path
from struct import unpack

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

import figure_paths

FONT = "Calibri"
INK = RGBColor(0x26, 0x32, 0x3A)
GRAY = RGBColor(0x3D, 0x47, 0x4E)

CAPTION_PT = 18
CAPTION_H = 0.34       # inches reserved under a panel for a one-line description
MARGIN = 0.12


def panels_dir(explicit=None) -> Path:
    """Where `prepare_overlay_panels.py` writes the bare panels these decks place."""
    if explicit:
        return Path(explicit).expanduser()
    return figure_paths.figures_out("overlay_panels")


def panel_path(name, panels=None) -> Path:
    return panels_dir(panels) / f"{name}.jpg"


def panel_aspect(name, panels=None) -> float:
    """height / width of a panel image, without pulling in an imaging library."""
    path = panel_path(name, panels)
    if not path.is_file():
        raise SystemExit(f"missing panel {path}; run prepare_overlay_panels.py first")
    data = path.read_bytes()
    # The JPEG SOF0/SOF2 marker carries the dimensions two bytes in.
    offset = 2
    while offset < len(data):
        marker, length = unpack(">HH", data[offset:offset + 4])
        if marker in (0xFFC0, 0xFFC2):
            height, width = unpack(">HH", data[offset + 5:offset + 9])
            return height / width
        offset += 2 + length
    raise SystemExit(f"no SOF marker in {path}")


def new_deck(width_in, height_in):
    """A one-slide deck of exactly this size, on the blank layout."""
    deck = Presentation()
    deck.slide_width = Inches(width_in)
    deck.slide_height = Inches(height_in)
    return deck, deck.slides.add_slide(deck.slide_layouts[6])


def add_text(slide, text, left, top, width, height, size, align=PP_ALIGN.CENTER,
             color=INK, bold=False, anchor=MSO_ANCHOR.TOP, vertical=False):
    """One text box, with the margins zeroed so the text sits where it is placed."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = anchor
    frame.word_wrap = True
    if vertical:
        # Bottom-to-top text. python-pptx has no setter for it, so write the attribute
        # the way PowerPoint and LibreOffice both read it.
        frame._txBody.bodyPr.set("vert", "vert270")
    run = frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.name = FONT
    run.font.color.rgb = color
    run.font.bold = bold
    frame.paragraphs[0].alignment = align
    return box


def save(deck, out_dir, name, width_in, height_in) -> Path:
    """Write the deck and say where it went and how big the slide is."""
    out_dir = Path(out_dir).expanduser()
    out_dir.mkdir(parents=True, exist_ok=True)
    path = out_dir / f"{name}.pptx"
    deck.save(str(path))
    print(f"wrote {path}  ({width_in:.2f} x {height_in:.2f} in)")
    return path
